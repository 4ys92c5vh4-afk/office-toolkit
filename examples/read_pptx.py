#!/usr/bin/env python3
"""
读取 PowerPoint 文件 (.pptx)
支持：文本提取、幻灯片结构、备注读取
"""

import sys
from pathlib import Path


def read_pptx_with_markitdown(filepath):
    """使用 markitdown 提取文本"""
    import subprocess
    try:
        result = subprocess.run(
            ["python3", "-m", "markitdown", filepath],
            capture_output=True, text=True, check=True
        )
        return result.stdout
    except Exception as e:
        print(f"⚠️ markitdown 失败: {e}")
        return None


def read_pptx_with_python(filepath):
    """使用 python-pptx 详细读取"""
    from pptx import Presentation

    prs = Presentation(filepath)
    print(f"📊 幻灯片数量: {len(prs.slides)}")

    all_text = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text = [f"\n{'='*60}\n📑 第 {i} 页\n{'='*60}"]

        # 提取所有形状中的文本
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())

            # 表格
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    row_text = [cell.text for cell in row.cells]
                    rows.append(" | ".join(row_text))
                slide_text.append("\n[表格]\n" + "\n".join(rows))

        # 备注
        if slide.has_notes_slide and slide.notes_slide.text_frame.text.strip():
            slide_text.append(f"\n💬 备注: {slide.notes_slide.text_frame.text.strip()}")

        all_text.append("\n".join(slide_text))

    return "\n".join(all_text)


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法: python3 read_pptx.py <文件路径>")
        sys.exit(1)

    filepath = sys.argv[1]
    if not Path(filepath).exists():
        print(f"❌ 文件不存在: {filepath}")
        sys.exit(1)

    print(f"🎬 读取 PPT 文件: {filepath}")

    # 优先使用 python-pptx（信息更详细）
    text = read_pptx_with_python(filepath)
    if text:
        print(text)

    # 备选：markitdown
    md_text = read_pptx_with_markitdown(filepath)
    if md_text:
        print("\n\n📋 Markdown 格式输出:")
        print(md_text[:2000])
